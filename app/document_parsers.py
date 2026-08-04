"""Document parsers for PDF, DOCX, XLSX, PPTX."""
from io import BytesIO
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation


def parse_pdf(file) -> dict:
    text = []
    tables = []
    try:
        with pdfplumber.open(file) as pdf:
            for p in pdf.pages:
                text.append(p.extract_text() or "")
                # extract tables
                try:
                    for table in p.extract_tables():
                        df = pd.DataFrame(table[1:], columns=table[0]) if len(table) > 1 else pd.DataFrame(table)
                        tables.append(df)
                except Exception:
                    pass
    except Exception:
        # file may be a BytesIO
        file.seek(0)
        with pdfplumber.open(file) as pdf:
            for p in pdf.pages:
                text.append(p.extract_text() or "")
    return {"text": "\n".join(text), "tables": tables}


def parse_docx(file) -> dict:
    text = []
    tables = []
    doc = Document(file)
    for p in doc.paragraphs:
        text.append(p.text)
    for t in doc.tables:
        rows = []
        for r in t.rows:
            rows.append([c.text for c in r.cells])
        if rows:
            df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows) > 1 else pd.DataFrame(rows)
            tables.append(df)
    return {"text": "\n".join(text), "tables": tables}


def parse_xlsx(file) -> dict:
    # read all sheets and return as tables
    try:
        xls = pd.read_excel(file, sheet_name=None)
    except Exception:
        file.seek(0)
        xls = pd.read_excel(file, sheet_name=None)
    tables = []
    for name, df in xls.items():
        tables.append(df)
    text = "\n".join([f"Sheet: {n}" for n in xls.keys()])
    return {"text": text, "tables": tables}


def parse_pptx(file) -> dict:
    prs = Presentation(file)
    text = []
    tables = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.has_table:
                table = shape.table
                rows = []
                for r in table.rows:
                    rows.append([c.text for c in r.cells])
                if rows:
                    df = pd.DataFrame(rows[1:], columns=rows[0]) if len(rows) > 1 else pd.DataFrame(rows)
                    tables.append(df)
        text.append("\n".join(slide_text))
    return {"text": "\n\n".join(text), "tables": tables}


def parse_file(uploaded_file) -> dict:
    # uploaded_file is Streamlit UploadedFile which behaves like BytesIO
    name = uploaded_file.name.lower()
    if name.endswith(".pdf"):
        return parse_pdf(uploaded_file)
    if name.endswith(".docx") or name.endswith(".doc"):
        return parse_docx(uploaded_file)
    if name.endswith(".xlsx") or name.endswith(".xls"):
        return parse_xlsx(uploaded_file)
    if name.endswith(".pptx") or name.endswith(".ppt"):
        return parse_pptx(uploaded_file)
    # fallback: try reading as text
    try:
        uploaded_file.seek(0)
        return {"text": uploaded_file.read().decode("utf-8", errors="ignore"), "tables": []}
    except Exception:
        return {"text": "", "tables": []}
